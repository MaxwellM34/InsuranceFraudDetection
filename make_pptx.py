"""Generate the Alan fraud-detection presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x24, 0x40)
RED    = RGBColor(0xD6, 0x28, 0x39)
ORANGE = RGBColor(0xE0, 0x7B, 0x39)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
GOLD   = RGBColor(0xF5, 0xA6, 0x23)
BLUE   = RGBColor(0x1A, 0x78, 0xC2)
TEAL   = RGBColor(0x0D, 0x7A, 0x8E)
LRED   = RGBColor(0xFF, 0xEB, 0xEB)
LORG   = RGBColor(0xFF, 0xF3, 0xE6)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Primitives ────────────────────────────────────────────────────────────────

def rect(sl, x, y, w, h, fill=None, line=None):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def tb(sl, text, x, y, w, h, size=13, bold=False, color=DGRAY,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    t = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    t.word_wrap = wrap
    tf = t.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return t


def header(sl, title, subtitle=None, q_tag=None):
    rect(sl, 0, 0, 13.33, 1.05, fill=NAVY)
    if q_tag:
        rect(sl, 0.35, 0.13, 1.2, 0.32, fill=GOLD)
        tb(sl, q_tag, 0.35, 0.13, 1.2, 0.32,
           size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(sl, title, 1.7, 0.08, 11.3, 0.65, size=24, bold=True, color=WHITE)
    else:
        tb(sl, title, 0.4, 0.1, 12.5, 0.65, size=24, bold=True, color=WHITE)
    if subtitle:
        rect(sl, 0, 1.05, 13.33, 0.38, fill=LGRAY)
        tb(sl, subtitle, 0.4, 1.08, 12.5, 0.33, size=12, color=DGRAY, italic=True)


def footer(sl):
    rect(sl, 0, 7.12, 13.33, 0.38, fill=NAVY)
    tb(sl, "Alan Assurance · Fraud Detection Case · 2024",
       0.3, 7.13, 12.7, 0.36, size=9, color=WHITE, align=PP_ALIGN.CENTER)


def score_pill(sl, score, x, y, w=1.55, h=0.52):
    c = RED if score >= 71 else (ORANGE if score >= 1 else GREEN)
    rect(sl, x, y, w, h, fill=c)
    tb(sl, f"{score}/100", x, y, w, h, size=17, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)


def status_pill(sl, status, x, y):
    cm = {
        "auto_held":     (RED,    "PAYMENTS HELD"),
        "needs_review":  (ORANGE, "NEEDS REVIEW"),
        "auto_approved": (GREEN,  "APPROVED"),
    }
    c, lbl = cm.get(status, (DGRAY, status.upper()))
    rect(sl, x, y, 1.85, 0.38, fill=c)
    tb(sl, lbl, x, y, 1.85, 0.38, size=9, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)


def section_divider(sl, q_num, q_label, title, description):
    rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(sl, 0, 5.8,  13.33, 1.7,  fill=RED)
    rect(sl, 0, 0,    0.18,  7.5,  fill=GOLD)
    rect(sl, 0.5, 1.6, 2.2, 0.55,  fill=GOLD)
    tb(sl, f"QUESTION {q_num}", 0.5, 1.6, 2.2, 0.55,
       size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tb(sl, q_label,     0.5, 2.3,  12.5, 0.5,  size=16, color=MGRAY, italic=True)
    tb(sl, title,       0.5, 2.95, 12.5, 1.5,  size=40, bold=True, color=WHITE)
    tb(sl, description, 0.5, 5.95, 12.5, 0.6,  size=13, color=WHITE, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)

rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 5.3, 13.33, 2.2, fill=RED)
rect(sl, 0, 0, 0.3, 7.5, fill=GOLD)

tb(sl, "ALAN ASSURANCE", 0.8, 0.9, 11.7, 0.55,
   size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
tb(sl, "Ops Runner — Technical Case", 0.8, 1.55, 11.7, 1.0,
   size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Detecting Fraud in Optical Care Claims", 0.8, 2.65, 11.7, 0.65,
   size=24, color=MGRAY, align=PP_ALIGN.CENTER)

tb(sl, (
    "Alan reimburses optical care claims submitted by opticians. "
    "This case asks: can we detect which providers are billing fraudulently "
    "— and build a system to catch it automatically?"
), 0.9, 3.45, 11.5, 0.85, size=14, color=MGRAY, align=PP_ALIGN.CENTER, wrap=True)

tb(sl, "March 2024", 0.8, 6.35, 11.7, 0.45,
   size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Dataset Overview
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "The Dataset — What We're Working With")
footer(sl)

# Metric cards
for val, label, sub, color, cx in [
    ("12",   "Optical providers",    "Independent optician shops",  NAVY,   0.35),
    ("221",  "Reimbursement claims", "Records of money paid out",   BLUE,   3.5),
    ("18",   "Months of history",    "January 2022 – June 2023",    TEAL,   6.65),
    ("2",    "Claim categories",     "Lunettes (glasses) · Lentilles (contacts)", ORANGE, 9.8),
]:
    rect(sl, cx+0.04, 1.62, 3.0, 1.55, fill=RGBColor(0xE4, 0xE4, 0xE4))
    rect(sl, cx, 1.58, 3.0, 1.55, fill=WHITE, line=MGRAY)
    rect(sl, cx, 1.58, 3.0, 0.2, fill=color)
    tb(sl, val,   cx+0.15, 1.86, 2.7, 0.62, size=38, bold=True, color=color)
    tb(sl, label, cx+0.15, 2.5,  2.7, 0.35, size=12, bold=True, color=NAVY)
    tb(sl, sub,   cx+0.15, 2.84, 2.7, 0.28, size=10, color=DGRAY, italic=True, wrap=True)

# Left panel — what a row looks like
rect(sl, 0.35, 3.52, 5.95, 3.15, fill=LGRAY)
rect(sl, 0.35, 3.52, 0.12, 3.15, fill=NAVY)
tb(sl, "What a single data row looks like",
   0.62, 3.6, 5.5, 0.38, size=13, bold=True, color=NAVY)
tb(sl, "Each row = one provider's total billings for one category in one month",
   0.62, 3.98, 5.5, 0.35, size=11, color=DGRAY, italic=True)

for j, (col, ex) in enumerate([
    ("Provider",  "Kylian's Frames"),
    ("Category",  "Lunettes (glasses)"),
    ("Month",     "March 2023"),
    ("Amount",    "€ 1 150.00  billed to Alan"),
]):
    bg = WHITE if j % 2 == 0 else LGRAY
    rect(sl, 0.55, 4.4 + j*0.4, 5.55, 0.38, fill=bg)
    tb(sl, col, 0.68, 4.42+j*0.4, 2.0, 0.28, size=11, bold=True, color=NAVY)
    tb(sl, ex,  2.85, 4.42+j*0.4, 3.1, 0.28, size=11, color=DGRAY)

# Right panel — the question
rect(sl, 6.6, 3.52, 6.4, 3.15, fill=WHITE, line=MGRAY)
rect(sl, 6.6, 3.52, 0.12, 3.15, fill=RED)
tb(sl, "The question we're trying to answer",
   6.88, 3.6, 5.9, 0.38, size=13, bold=True, color=NAVY)

for j, (q, a) in enumerate([
    ("Who might be committing fraud?",
     "Some opticians may be billing for services that were never provided — or inflating bills to extract more money from Alan."),
    ("What does fraud look like in this data?",
     "Unusual billing patterns: sudden spikes, identical repeated amounts, or always claiming both glasses and contacts together."),
    ("What can we do about it?",
     "Automatically score each provider and surface the suspicious ones for a human reviewer to investigate."),
]):
    rect(sl, 6.75, 4.05+j*0.82, 6.1, 0.75, fill=LGRAY if j % 2 == 0 else WHITE)
    tb(sl, q, 6.9, 4.08+j*0.82, 5.8, 0.3, size=10, bold=True, color=NAVY, wrap=True)
    tb(sl, a, 6.9, 4.38+j*0.82, 5.8, 0.38, size=10, color=DGRAY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Q1 Divider
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
section_divider(sl,
    q_num=1,
    q_label="Data Exploration",
    title="What suspicious\npatterns did we find?",
    description="8 out of 12 providers showed at least one red flag · 4 distinct fraud patterns identified")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Results Overview
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Results — Which Providers Look Suspicious?", q_tag="Q1",
       subtitle="Each provider is given a risk score 0–100 based on how many suspicious patterns were detected")
footer(sl)

providers = [
    ("Queen optics",         99, "auto_held",     "Co-billing + billing spike + round numbers"),
    ("Mike lunettes",        96, "auto_held",     "Co-billing + billing spike + round numbers"),
    ("Les lunettes à Soso",  84, "auto_held",     "Co-billing + spike + repeated + round numbers"),
    ("Penthievre alambics",  78, "auto_held",     "Co-billing + billing spike + round numbers"),
    ("Runner glasses",       74, "auto_held",     "Co-billing + repeated amounts + round numbers"),
    ("Kylian's frames",      54, "needs_review",  "Co-billing + billing spike + round numbers"),
    ("Roudoudou lentilles",  40, "needs_review",  "Co-billing + round numbers"),
    ("Voodoo",               15, "needs_review",  "Round number billing detected"),
    ("La classe à Dallas",    0, "auto_approved", "No suspicious patterns found"),
    ("Cool optics",           0, "auto_approved", "No suspicious patterns found"),
    ("abc optic",             0, "auto_approved", "No suspicious patterns found"),
    ("22 optics",             0, "auto_approved", "No suspicious patterns found"),
]

xs = [0.3, 4.05, 5.65, 7.55]
rect(sl, 0.3, 1.55, 12.73, 0.4, fill=NAVY)
for lbl, x in zip(["Provider", "Risk Score", "Outcome", "What was detected"], xs):
    tb(sl, lbl, x+0.1, 1.57, 1.9, 0.34, size=10, bold=True, color=WHITE)

ry = 1.98
for name, score, status, flags in providers:
    row_bg = {
        "auto_held":     RGBColor(0xFF, 0xED, 0xED),
        "needs_review":  RGBColor(0xFF, 0xF6, 0xED),
        "auto_approved": RGBColor(0xF0, 0xF9, 0xF2),
    }.get(status, LGRAY)
    rect(sl, 0.3, ry, 12.73, 0.39, fill=row_bg)
    sc = RED if score >= 71 else (ORANGE if score >= 1 else GREEN)
    outcome = {"auto_held": "Payments held", "needs_review": "Under review",
               "auto_approved": "Cleared"}[status]
    tb(sl, name,    xs[0]+0.1, ry+0.05, 3.6, 0.28, size=11, color=DGRAY)
    tb(sl, str(score), xs[1]+0.1, ry+0.05, 1.4, 0.28, size=12, bold=True, color=sc)
    tb(sl, outcome, xs[2]+0.1, ry+0.05, 1.75, 0.28, size=10, bold=True, color=sc)
    tb(sl, flags,   xs[3]+0.1, ry+0.05, 5.4,  0.28, size=10, color=DGRAY)
    ry += 0.39

rect(sl, 0.3, ry+0.08, 12.73, 0.44, fill=NAVY)
tb(sl, "5 providers had payments automatically held  ·  3 flagged for manual review  ·  4 cleared",
   0.5, ry+0.12, 12.3, 0.34, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Pattern 1: Co-billing
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Pattern 1 — Glasses and Contacts Always Billed Together", q_tag="Q1",
       subtitle="Some opticians systematically claim both Lunettes and Lentilles for the same month, month after month")
footer(sl)

# What's suspicious box
rect(sl, 0.3, 1.55, 12.73, 1.0, fill=RGBColor(0xFF, 0xF3, 0xE6))
rect(sl, 0.3, 1.55, 0.12, 1.0, fill=ORANGE)
tb(sl, "Why is this suspicious?", 0.55, 1.62, 5.0, 0.32, size=12, bold=True, color=NAVY)
tb(sl, (
    "In reality, a single patient can't wear both glasses and contact lenses at the same time. "
    "An optician billing for both products every single month is likely adding fictitious contact-lens "
    "charges to genuine glasses sales — the patient pays less out-of-pocket while Alan absorbs the extra cost."
), 0.55, 1.94, 12.1, 0.55, size=11, color=DGRAY, wrap=True)

# two provider cards
for xi, name, score, status, pct, n_active, n_dual, extra in [
    (0.3,  "Kylian's Frames",    54, "needs_review",
     76, 13, 10, "Also had one month where billings suddenly spiked 5× above normal"),
    (6.68, "Roudoudou Lentilles", 40, "needs_review",
     100, 12, 12, "Every single active month had both categories — 100% of the time"),
]:
    rect(sl, xi, 2.68, 6.28, 4.12, fill=WHITE, line=MGRAY)
    rect(sl, xi, 2.68, 6.28, 0.08, fill=ORANGE)
    tb(sl, name, xi+0.2, 2.84, 4.5, 0.38, size=14, bold=True, color=NAVY)
    score_pill(sl, score, xi+4.65, 2.79, w=1.55, h=0.44)

    # big stat
    rect(sl, xi+0.2, 3.35, 5.85, 0.75, fill=LORG)
    tb(sl, f"{pct}% of active months", xi+0.3, 3.38, 5.6, 0.42,
       size=22, bold=True, color=ORANGE)
    tb(sl, f"had both Lunettes and Lentilles billed  ({n_dual} out of {n_active} months)",
       xi+0.3, 3.8, 5.6, 0.28, size=11, color=DGRAY)

    tb(sl, extra, xi+0.2, 4.22, 5.85, 0.38, size=11, color=DGRAY, italic=True, wrap=True)
    status_pill(sl, status, xi+0.2, 4.72)
    sc = ORANGE
    tb(sl, f"Risk score: {score}/100", xi+2.25, 4.76, 4.0, 0.32, size=12, bold=True, color=sc)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Pattern 2: Billing Spikes + Multi-Rule Providers
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Pattern 2 — Sudden Unexplained Billing Spikes", q_tag="Q1",
       subtitle="Some providers had months where billings jumped dramatically above their own historical average")
footer(sl)

# What's suspicious
rect(sl, 0.3, 1.55, 12.73, 0.78, fill=LRED)
rect(sl, 0.3, 1.55, 0.12,  0.78, fill=RED)
tb(sl, "Why is this suspicious?", 0.55, 1.62, 4.0, 0.32, size=12, bold=True, color=NAVY)
tb(sl, (
    "A sudden spike well above a provider's own normal billing level can indicate a burst of "
    "fictitious or inflated claims — especially when it coincides with the co-billing pattern above."
), 0.55, 1.93, 12.1, 0.35, size=11, color=DGRAY, wrap=True)

# three provider cards
spike_providers = [
    ("Mike Lunettes",        96, "auto_held", 3, "85% co-billing rate across all months"),
    ("Penthievre Alambics",  78, "auto_held", 2, "94% co-billing rate — highest in dataset"),
    ("Queen Optics",         99, "auto_held", 3, "96% co-billing rate — near-universal"),
]
for k, (name, score, status, n_spikes, cobill_note) in enumerate(spike_providers):
    xi = 0.3 + k * 4.32
    rect(sl, xi, 2.45, 4.1, 4.3, fill=WHITE, line=MGRAY)
    rect(sl, xi, 2.45, 4.1, 0.08, fill=RED)
    tb(sl, name, xi+0.15, 2.62, 3.55, 0.38, size=13, bold=True, color=NAVY)
    score_pill(sl, score, xi+2.55, 2.57, w=1.45, h=0.44)
    status_pill(sl, status, xi+0.15, 3.12)

    rect(sl, xi+0.15, 3.63, 3.8, 0.62, fill=LRED)
    tb(sl, f"{n_spikes} spike month{'s' if n_spikes > 1 else ''}",
       xi+0.25, 3.66, 3.55, 0.38, size=18, bold=True, color=RED)
    tb(sl, "where billings were > 5× their own 6-month average",
       xi+0.25, 4.0, 3.55, 0.28, size=10, color=DGRAY, wrap=True)

    rect(sl, xi+0.15, 4.4, 3.8, 0.7, fill=LGRAY)
    tb(sl, "Also: co-billing", xi+0.25, 4.45, 3.55, 0.28, size=10, bold=True, color=NAVY)
    tb(sl, cobill_note, xi+0.25, 4.73, 3.55, 0.35, size=10, color=DGRAY, wrap=True)

    tb(sl, "Two independent signals → score 100", xi+0.15, 5.28, 3.8, 0.32,
       size=10, bold=True, color=RED, wrap=True)

tb(sl, "When two unrelated fraud signals appear simultaneously it dramatically increases confidence the provider is acting fraudulently.",
   0.3, 6.65, 12.7, 0.38, size=11, color=DGRAY, italic=True, align=PP_ALIGN.CENTER, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Pattern 3: Repeated Amounts + Most Complex Case
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Pattern 3 — Identical Amounts Billed Repeatedly", q_tag="Q1",
       subtitle="The exact same euro amount appearing 3 or more times within a 12-month period")
footer(sl)

# What's suspicious
rect(sl, 0.3, 1.55, 12.73, 0.75, fill=LGRAY)
rect(sl, 0.3, 1.55, 0.12,  0.75, fill=BLUE)
tb(sl, "Why is this suspicious?", 0.55, 1.62, 4.0, 0.32, size=12, bold=True, color=NAVY)
tb(sl, (
    "Real claims vary naturally. When the exact same euro amount appears over and over for the same "
    "provider and product category, it suggests the claims may be copy-pasted fabrications rather than genuine billing."
), 0.55, 1.93, 12.1, 0.32, size=11, color=DGRAY, wrap=True)

# Runner Glasses
rect(sl, 0.3, 2.42, 5.95, 4.38, fill=WHITE, line=MGRAY)
rect(sl, 0.3, 2.42, 5.95, 0.08, fill=BLUE)
tb(sl, "Runner Glasses", 0.5, 2.58, 3.6, 0.38, size=14, bold=True, color=NAVY)
score_pill(sl, 74, 4.55, 2.53, w=1.28, h=0.44)
status_pill(sl, "auto_held", 0.5, 3.04)

tb(sl, "Three separate repeated-amount patterns detected:", 0.5, 3.55, 5.5, 0.32,
   size=11, bold=True, color=NAVY)

rect(sl, 0.5, 3.9, 5.55, 0.35, fill=NAVY)
for lbl, xi in zip(["Amount", "Category", "How many times"], [0.5, 2.1, 3.7]):
    tb(sl, lbl, xi+0.1, 3.92, 1.5, 0.27, size=9, bold=True, color=WHITE)

for j, (amt, cat, occ) in enumerate([
    ("€ 850.00", "Glasses",  "4 times in 12 months"),
    ("€ 320.00", "Contacts", "3 times in 12 months"),
    ("€ 650.00", "Glasses",  "3 times in 12 months"),
]):
    bg = LGRAY if j % 2 == 0 else WHITE
    rect(sl, 0.5, 4.28+j*0.38, 5.55, 0.36, fill=bg)
    tb(sl, amt, 0.62, 4.3+j*0.38, 1.5, 0.26, size=11, bold=True, color=NAVY)
    tb(sl, cat, 2.22, 4.3+j*0.38, 1.5, 0.26, size=11, color=DGRAY)
    tb(sl, occ, 3.82, 4.3+j*0.38, 2.1, 0.26, size=11, color=DGRAY)

tb(sl, "Also had co-billing in 55% of months → total score 74",
   0.5, 5.5, 5.5, 0.32, size=11, bold=True, color=RED)

# Les lunettes à Soso — all 3 patterns
rect(sl, 6.65, 2.42, 6.35, 4.38, fill=WHITE, line=MGRAY)
rect(sl, 6.65, 2.42, 6.35, 0.08, fill=RED)
tb(sl, "Les Lunettes à Soso", 6.85, 2.58, 4.0, 0.38, size=14, bold=True, color=NAVY)
score_pill(sl, 84, 11.4, 2.53, w=1.28, h=0.44)
status_pill(sl, "auto_held", 6.85, 3.04)

rect(sl, 6.85, 3.55, 6.0, 0.4, fill=RED)
tb(sl, "The only provider to trigger all three patterns",
   6.95, 3.58, 5.8, 0.32, size=11, bold=True, color=WHITE)

for j, (label, detail) in enumerate([
    ("Co-billing",          "77% of months had both glasses + contacts billed"),
    ("Billing spike",       "1 month jumped to over 5× their normal billing level"),
    ("Repeated amounts",    "2 separate amounts repeated 3+ times each"),
]):
    rect(sl, 6.85, 4.0+j*0.62, 6.0, 0.55, fill=LGRAY if j % 2 == 0 else WHITE)
    tb(sl, label,  7.0, 4.03+j*0.62, 2.2, 0.32, size=11, bold=True, color=NAVY)
    tb(sl, detail, 9.3, 4.03+j*0.62, 3.4, 0.45, size=10, color=DGRAY, wrap=True)

rect(sl, 6.85, 5.9, 6.0, 0.75, fill=LRED)
tb(sl, "Three separate, independent signals → highest confidence fraud case in the dataset.",
   7.0, 5.97, 5.7, 0.6, size=11, color=RED, bold=True, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Pattern 4: Round Number Billing
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Pattern 4 — Systematic Round-Number Billing", q_tag="Q1",
       subtitle="Providers who consistently bill exact round sums (e.g. 300€, 600€) rather than realistic irregular amounts")
footer(sl)

# What's suspicious box
rect(sl, 0.3, 1.55, 12.73, 0.85, fill=RGBColor(0xE8, 0xF4, 0xFD))
rect(sl, 0.3, 1.55, 0.12, 0.85, fill=BLUE)
tb(sl, "Why is this suspicious?", 0.55, 1.62, 5.0, 0.32, size=12, bold=True, color=NAVY)
tb(sl, (
    "Real optical purchases — glasses frames, lenses, fitting — never add up to a perfectly round number. "
    "Amounts like €349.91 or €127.50 reflect actual product pricing. When a provider bills €300.00, €600.00, "
    "€200.00 every single month with no variation in cents, it strongly suggests fabricated charges rather than genuine sales."
), 0.55, 1.93, 12.1, 0.42, size=11, color=DGRAY, wrap=True)

# Runner Glasses — main example
rect(sl, 0.3, 2.55, 6.1, 4.2, fill=WHITE, line=MGRAY)
rect(sl, 0.3, 2.55, 6.1, 0.08, fill=BLUE)
tb(sl, "Runner Glasses — Contacts billing", 0.5, 2.7, 5.5, 0.38, size=13, bold=True, color=NAVY)

rect(sl, 0.5, 3.18, 5.7, 0.62, fill=LRED)
tb(sl, "100% of Lentilles claims are round multiples of 100€",
   0.62, 3.21, 5.45, 0.42, size=13, bold=True, color=RED)

rect(sl, 0.5, 3.92, 5.7, 0.35, fill=NAVY)
for lbl, xi2 in zip(["Period", "Amount", "Round?"], [0.5, 2.8, 4.7]):
    tb(sl, lbl, xi2+0.1, 3.94, 1.8, 0.27, size=9, bold=True, color=WHITE)

for j, (period, amt, is_round) in enumerate([
    ("April 2022",     "300,00 €", True),
    ("June 2022",      "300,00 €", True),
    ("September 2022", "300,00 €", True),
    ("October 2022",   "600,00 €", True),
    ("December 2022",  "600,00 €", True),
    ("April 2023",     "200,00 €", True),
]):
    bg = LGRAY if j % 2 == 0 else WHITE
    rect(sl, 0.5, 4.3+j*0.34, 5.7, 0.32, fill=bg)
    tb(sl, period, 0.62, 4.32+j*0.34, 2.1, 0.24, size=10, color=DGRAY)
    tb(sl, amt,    2.92, 4.32+j*0.34, 1.7, 0.24, size=10, bold=True, color=RED if is_round else DGRAY)
    tb(sl, "✓ Round" if is_round else "—", 4.82, 4.32+j*0.34, 1.2, 0.24,
       size=9, bold=is_round, color=RED if is_round else DGRAY)

# Right panel — comparison + explanation
rect(sl, 6.6, 2.55, 6.4, 4.2, fill=WHITE, line=MGRAY)
rect(sl, 6.6, 2.55, 6.4, 0.08, fill=NAVY)
tb(sl, "Compare: Lunettes (glasses) — same provider", 6.8, 2.7, 6.0, 0.38, size=12, bold=True, color=NAVY)
tb(sl, "Glasses amounts look realistic — irregular cents, normal price variation",
   6.8, 3.12, 5.95, 0.32, size=10, color=DGRAY, italic=True)

for j2, (period, amt) in enumerate([
    ("April 2022",     "349,91 €"),
    ("June 2022",      "449,91 €"),
    ("September 2022", "349,91 €"),
    ("October 2022",   "699,82 €"),
    ("December 2022",  "349,91 €"),
    ("April 2023",     "564,82 €"),
]):
    bg = LGRAY if j2 % 2 == 0 else WHITE
    rect(sl, 6.8, 3.55+j2*0.34, 5.95, 0.32, fill=bg)
    tb(sl, period, 6.92, 3.57+j2*0.34, 2.1, 0.24, size=10, color=DGRAY)
    tb(sl, amt,    9.12, 3.57+j2*0.34, 1.7, 0.24, size=10, bold=True, color=GREEN)
    tb(sl, "Realistic", 10.92, 3.57+j2*0.34, 1.2, 0.24, size=9, color=GREEN)

rect(sl, 6.8, 5.65, 5.95, 0.75, fill=LRED)
tb(sl, (
    "Same provider, same months — glasses have real prices, contacts are always "
    "round numbers. The contacts charges appear fabricated."
), 6.92, 5.7, 5.7, 0.65, size=10, bold=True, color=RED, wrap=True)

tb(sl, "This pattern also appears in Roudoudou Lentilles, Kylian's Frames, and most other flagged providers.",
   0.3, 6.85, 12.7, 0.32, size=10, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Q2 Divider
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
section_divider(sl,
    q_num=2,
    q_label="Detection System Design",
    title="How did we turn these\npatterns into a system?",
    description="4 automated detection rules · risk scoring pipeline · instant routing of every provider")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Detection Engine
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "The Detection Engine — 4 Rules, Automatically Applied to Every Provider",
       q_tag="Q2",
       subtitle="Each rule adds points to a provider's risk score. Final score 0–100 determines what happens next.")
footer(sl)

def rule_card(sl, num, title, plain_english, logic_line, score_line, x, w=3.05):
    rect(sl, x, 1.55, w, 3.85, fill=NAVY)
    rect(sl, x, 1.55, w, 0.48, fill=RED)
    tb(sl, f"Rule {num}", x+0.15, 1.58, w-0.3, 0.38, size=14, bold=True, color=WHITE)
    tb(sl, title, x+0.15, 2.1, w-0.3, 0.4, size=12, bold=True, color=GOLD)
    tb(sl, plain_english, x+0.15, 2.56, w-0.3, 0.78, size=9, color=MGRAY, wrap=True)
    rect(sl, x, 3.38, w, 0.3, fill=RGBColor(0x26, 0x35, 0x5C))
    tb(sl, "Logic:", x+0.15, 3.4, 0.6, 0.25, size=8, bold=True, color=GOLD)
    tb(sl, logic_line, x+0.75, 3.4, w-0.9, 0.25, size=8, color=WHITE)
    rect(sl, x, 5.0, w, 0.42, fill=GOLD)
    tb(sl, score_line, x+0.08, 5.02, w-0.16, 0.36,
       size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

rule_card(sl, 1, "Billing Spike",
    "Flag any month where a provider's total billings are more than 5× higher than "
    "their own rolling 6-month average for that product category.",
    "current month > 5 × 6-month median",
    "+20 points per spiked month",
    0.3)

rule_card(sl, 2, "Systematic Co-Billing",
    "Flag providers who bill both glasses and contacts in the same month for the "
    "majority of their active months (≥ 50%). Occasional overlap is normal; "
    "systematic overlap is not.",
    "≥ 50% of months have both categories",
    "12–25 pts  (scales with rate)",
    3.52)

rule_card(sl, 3, "Repeated Identical Amount",
    "Flag any (provider, category) pair where the exact same euro amount appears "
    "3 or more times within any rolling 12-month window. Real billing varies naturally.",
    "same € amount ≥ 3× in 12-month window",
    "+15 points per pattern",
    6.74)

rule_card(sl, 4, "Round Number Billing",
    "Flag providers where ≥ 70% of claims in a category are whole euro amounts "
    "(no cents). Real optical purchases have irregular amounts like 349.91€.",
    "≥ 70% of claims are whole euros",
    "+15 points per category",
    9.96)

# Routing
rect(sl, 0.3, 5.55, 12.73, 0.35, fill=NAVY)
tb(sl, "What happens to the score", 0.4, 5.57, 12.5, 0.28,
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for lbl, rng, desc, c, xi in [
    ("Cleared",        "Score = 0",    "No flags — approved automatically",       GREEN,  0.3),
    ("Under Review",   "Score 1–70",   "At least one signal — human reviews it",  ORANGE, 4.55),
    ("Payments Held",  "Score > 70",   "Strong evidence — payments suspended",    RED,    8.8),
]:
    rect(sl, xi, 5.95, 3.85, 1.1, fill=c)
    tb(sl, lbl,  xi+0.1, 5.97, 3.65, 0.4,  size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, rng,  xi+0.1, 6.37, 3.65, 0.3,  size=11, color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, desc, xi+0.1, 6.67, 3.65, 0.32, size=9,  color=WHITE, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — App Introduction / Demo Transition
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)

rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
rect(sl, 0, 0, 0.25, 7.5, fill=GOLD)
rect(sl, 0, 6.05, 13.33, 1.45, fill=RGBColor(0x0D, 0x17, 0x30))

tb(sl, "Beyond the analysis —", 0.55, 0.7, 10.0, 0.52,
   size=20, color=MGRAY, italic=True)
tb(sl, "we built a working system\naround it", 0.55, 1.3, 11.5, 1.6,
   size=44, bold=True, color=WHITE)

tb(sl, (
    "The detection logic above is packaged into a full web application "
    "that an ops team can actually use day-to-day: upload new claims, "
    "run detection with one click, and review every flagged provider with full evidence."
), 0.55, 3.1, 9.5, 0.85, size=14, color=MGRAY, wrap=True)

# four feature chips
for i, (title_txt, desc, c) in enumerate([
    ("Dashboard",     "Live risk scores, flagged provider count, recent review activity",      BLUE),
    ("Provider list", "Run detection · search and filter · see every provider's risk score",   TEAL),
    ("Flag evidence", "Click any provider to see exactly which claims triggered each flag",    ORANGE),
    ("Review panel",  "Submit a decision on any provider — approve, escalate, or blacklist",   GREEN),
]):
    xi = 0.55 + i * 3.2
    rect(sl, xi, 4.1, 3.0, 1.7, fill=RGBColor(0x26, 0x35, 0x5C))
    rect(sl, xi, 4.1, 3.0, 0.1, fill=c)
    tb(sl, title_txt, xi+0.18, 4.28, 2.65, 0.38, size=13, bold=True, color=WHITE)
    tb(sl, desc,      xi+0.18, 4.68, 2.65, 0.9,  size=10, color=MGRAY, wrap=True)

tb(sl, "Built with  Next.js + TypeScript  ·  Python FastAPI  ·  PostgreSQL  ·  Full FR/EN localisation",
   0.55, 6.18, 12.2, 0.38, size=10,
   color=RGBColor(0x77, 0x88, 0xAA), align=PP_ALIGN.CENTER)

# webapp URL
rect(sl, 0.55, 5.62, 12.2, 0.38, fill=RGBColor(0x0D, 0x17, 0x30))
tb(sl, "🌐  Live demo:  alan-production-0d14.up.railway.app",
   0.65, 5.65, 12.0, 0.32, size=11, bold=True,
   color=GOLD, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Q3 Divider
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
section_divider(sl,
    q_num=3,
    q_label="Recommendations",
    title="What should happen\nwith these providers?",
    description="Immediate actions · system limitations · how to improve detection over time")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Payments Held (auto_held providers)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Immediate Action Required — Payments Held", q_tag="Q3",
       subtitle="5 providers automatically held — strong evidence across multiple independent fraud signals")
footer(sl)

held_providers = [
    ("Queen optics",         99, "Co-billing · billing spike · round numbers",
     "96% of months had both glasses + contacts. 3 spike months. All contacts claims are whole euros.",
     "Full audit + suspend all pending payments"),
    ("Mike lunettes",        96, "Co-billing · billing spike · round numbers",
     "85% co-billing rate. 3 spike months above 5× their 6-month average.",
     "Full audit + suspend all pending payments"),
    ("Les lunettes à Soso",  84, "Co-billing · spike · repeated amounts · round numbers",
     "Only provider to trigger all 4 rules. 4 completely independent fraud signals.",
     "Priority audit — highest confidence fraud case"),
    ("Penthievre alambics",  78, "Co-billing · billing spike · round numbers",
     "94% co-billing rate — highest in dataset. 2 spike months.",
     "Full audit + suspend all pending payments"),
    ("Runner glasses",       74, "Co-billing · repeated amounts · round numbers",
     "55% co-billing. 3 repeated-amount patterns. All contacts: exact round euros.",
     "Audit focus on fabricated contacts charges"),
]

sxs = [0.3, 3.85, 7.55]
rect(sl, 0.3, 1.55, 12.73, 0.35, fill=NAVY)
for lbl, x in zip(["Provider  /  Score", "Signals detected", "Recommended action"], sxs):
    tb(sl, lbl, x+0.12, 1.57, 3.4, 0.28, size=10, bold=True, color=WHITE)

for k, (name, score, signals, detail, action) in enumerate(held_providers):
    ry = 1.95 + k * 1.02
    rect(sl, 0.3, ry, 12.73, 0.98, fill=LRED)
    rect(sl, 0.3, ry, 0.1, 0.98, fill=RED)
    sc = RED
    tb(sl, name,  0.5,  ry+0.06, 3.2, 0.32, size=13, bold=True, color=NAVY)
    score_pill(sl, score, 0.5, ry+0.52, w=1.3, h=0.36)
    tb(sl, signals, 3.97, ry+0.06, 3.45, 0.28, size=10, bold=True, color=RED)
    tb(sl, detail,  3.97, ry+0.38, 3.45, 0.5,  size=9,  color=DGRAY, wrap=True)
    tb(sl, action,  7.67, ry+0.22, 5.2,  0.5,  size=10, bold=True, color=NAVY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Needs Review + Cleared
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "Under Review & Cleared Providers", q_tag="Q3",
       subtitle="3 providers flagged for manual review · 4 providers fully cleared with no suspicious patterns")
footer(sl)

# Under review section
rect(sl, 0.3, 1.55, 12.73, 0.35, fill=ORANGE)
tb(sl, "Flagged for manual review  —  score below auto-hold threshold but patterns present",
   0.45, 1.58, 12.3, 0.28, size=11, bold=True, color=WHITE)

review_providers = [
    ("Kylian's frames",     54, "Co-billing (76%) · billing spike · round numbers",
     "Below the 70-point hold threshold but has 3 signals. Recommend manual review before next payment cycle."),
    ("Roudoudou lentilles", 40, "Systematic co-billing (100% of months) · round numbers",
     "Every single active month had both categories billed. Low total due to no spike. Review co-billing pattern."),
    ("Voodoo",              15, "Round number billing detected",
     "Low score but whole-euro billing pattern present. Monitor — flag for review if pattern continues next cycle."),
]

for k, (name, score, signals, detail) in enumerate(review_providers):
    ry = 1.95 + k * 0.88
    rect(sl, 0.3, ry, 12.73, 0.82, fill=LORG)
    rect(sl, 0.3, ry, 0.1, 0.82, fill=ORANGE)
    score_pill(sl, score, 0.5, ry+0.22, w=1.3, h=0.36)
    tb(sl, name,    2.0,  ry+0.06, 3.5, 0.32, size=13, bold=True, color=NAVY)
    tb(sl, signals, 2.0,  ry+0.42, 3.5, 0.32, size=9,  bold=True, color=ORANGE)
    tb(sl, detail,  5.7,  ry+0.18, 7.1, 0.55, size=10, color=DGRAY, wrap=True)

# Cleared section
rect(sl, 0.3, 4.6, 12.73, 0.35, fill=GREEN)
tb(sl, "Cleared — no suspicious patterns detected",
   0.45, 4.63, 12.3, 0.28, size=11, bold=True, color=WHITE)

cleared = [("La classe à Dallas", 0), ("Cool optics", 0), ("abc optic", 0), ("22 optics", 0)]
for k, (name, score) in enumerate(cleared):
    cx = 0.3 + k * 3.18
    rect(sl, cx, 5.0, 3.08, 0.85, fill=RGBColor(0xF0, 0xF9, 0xF2), line=MGRAY)
    score_pill(sl, score, cx+0.12, 5.08, w=1.1, h=0.32)
    tb(sl, name, cx+1.3, 5.12, 1.7, 0.28, size=10, bold=True, color=GREEN, wrap=True)
    tb(sl, "No flags raised", cx+0.12, 5.5, 2.8, 0.28, size=9, color=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — How to Improve Over Time
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
header(sl, "How to Improve Detection Over Time", q_tag="Q3",
       subtitle="Current system catches clear patterns — here's how to make it more precise and harder to game")
footer(sl)

improvements = [
    (BLUE,   "Member-level data",
     "Get individual claim-level data instead of monthly aggregates.",
     "Currently we detect co-billing at the provider level — we can't confirm it's the same patient "
     "being billed for both glasses and contacts. With member IDs per claim, we could detect the "
     "exact fraud scenario: one patient billed for both products in the same month.",
     "Higher precision · fewer false positives · unambiguous evidence for legal action"),
    (TEAL,   "Geographic clustering",
     "Flag multiple providers operating from the same address.",
     "A network of opticians run by the same owner could each stay below the scoring threshold "
     "individually while collectively committing large-scale fraud. Address matching and network "
     "analysis would catch coordinated schemes invisible to per-provider rules.",
     "Detects organised fraud rings · cross-provider signal"),
    (ORANGE, "Seasonal calibration",
     "Adjust rolling medians to account for known seasonal spikes.",
     "Legitimate sales spikes occur every December (gift-giving) and September (back to school). "
     "A spike rule that doesn't account for seasonality will generate false positives for honest "
     "providers during these periods, wasting reviewer time.",
     "Fewer false positives · reviewer trust maintained"),
    (GREEN,  "Machine learning scoring",
     "Train a model on confirmed fraud cases to replace fixed point weights.",
     "The current 4-rule system uses fixed scores (+20, +25, +15, +15) that were chosen "
     "manually. A supervised ML model trained on confirmed fraud outcomes would learn which "
     "combinations of signals are most predictive — and weight them accordingly.",
     "Dynamic weights · adapts as fraud patterns evolve"),
]

for k, (c, title, subtitle_txt, body, benefit) in enumerate(improvements):
    row = k // 2
    col = k % 2
    bx = 0.3 + col * 6.45
    by = 1.55 + row * 2.62
    rect(sl, bx, by, 6.25, 2.5, fill=WHITE, line=MGRAY)
    rect(sl, bx, by, 6.25, 0.1, fill=c)
    tb(sl, title,        bx+0.2, by+0.18, 5.8, 0.35, size=14, bold=True, color=NAVY)
    tb(sl, subtitle_txt, bx+0.2, by+0.55, 5.8, 0.28, size=10, bold=True, color=c, italic=True)
    tb(sl, body,         bx+0.2, by+0.88, 5.8, 0.95, size=9,  color=DGRAY, wrap=True)
    rect(sl, bx+0.2, by+1.92, 5.8, 0.42, fill=LGRAY)
    tb(sl, f"→  {benefit}", bx+0.3, by+1.97, 5.6, 0.32, size=9, bold=True, color=c, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = "/home/maxwell/workspace/alan/interview/fraud_ops_presentation_updated.pptx"
prs.save(out)
print(f"Saved → {out}")
